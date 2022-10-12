from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

# for i in range(0, 11):
#     title_slide_layout = prs.slide_layouts[i] # 슬라이드 종류 선택
#     slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가

# prs.save('add all slides.pptx')

for i in range(0, 11):
    print("--------[%d] ------ "%(i))
    slide = prs.slides.add_slide(prs.slide_layouts[i])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))