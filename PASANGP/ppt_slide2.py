from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

bullet_slide_layout = prs.slide_layouts[1] # 0 : 제목 및 내용 슬라이드에 해당
slide = prs.slides.add_slide(bullet_slide_layout) # 슬라이드 추가

# 제목 - 제목에 값넣기
title = slide.placeholders[0] # 제목
title.text = "Bullet Slide" # 제목에 값 넣기

# 내용
body_content = slide.placeholders[1] # 제목 상자는 placeholders[0], 내용 상자는 [1]
tf = body_content.text_frame
tf.text = 'Fild the bullet slide layout'

#단락 추가
p = tf.add_paragraph()
p.text = 'Use _TextFrame.text for first bullet'
p.level = 1   # 1 : 들어쓰기 레벨

#단락 추가
p = tf.add_paragraph()
p.text = ' 서브 bullets 를 위한 _TextFrame.add_paragraph() 사용'
p.level = 2   # 2 : 들어쓰기 레벨

#단락 추가
p = tf.add_paragraph()
p.text = ' ABC 테스트'
p.level = 1   # 2 : 들어쓰기 레벨


#단락 추가
p = tf.add_paragraph()
p.text = ' 기아타이거즈 ABC 테스트'
p.level = 0   # 2 : 들어쓰기 레벨

# 카카오 토큰을 저장할 파일명
PPT_FILENAME = "res/powerpoint_handling/test2_2.pptx"


# 저장
prs.save(PPT_FILENAME)