import matplotlib.pyplot as plt
from pandas.plotting import table
import os
import pandas as pd
import requests
import datetime
from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표 등을 그리기 위해






# %matplotlib inline은 jupyter notebook 사용자 용 - jupyter notebook 내에 그래프가 그려지게 한다.
#%matplotlib inline

import smtplib

# 이메일 메시지에 다양한 형식을 중첩하여 담기 위한 객체
from email.mime.multipart import MIMEMultipart

# 이메일 메시지를 이진 데이터로 바꿔주는 인코더
from email import encoders

# 텍스트 형식
from email.mime.text import MIMEText
# 이미지 형식
from email.mime.image import MIMEImage
# 오디오 형식
from email.mime.audio import MIMEAudio

# 위의 모든 객체들을 생성할 수 있는 기본 객체
# MIMEBase(_maintype, _subtype)
# MIMEBase(<메인 타입>, <서브 타입>)
from email.mime.base import MIMEBase


def send_email(smtp_info, msg):
    with smtplib.SMTP(smtp_info["smtp_server"], smtp_info["smtp_port"]) as server:
        # TLS 보안 연결
        server.starttls()
        # 로그인
        server.login(smtp_info["smtp_user_id"], smtp_info["smtp_user_pw"])
        
        # 로그인 된 서버에 이메일 전송
        response = server.sendmail(msg['from'], msg['to'], msg.as_string()) # 메시지를 보낼때는 .as_string() 메소드를 사용해서 문자열로 바꿔줍니다.
        
        # 이메일을 성공적으로 보내면 결과는 {}
        if not response:
            print('이메일을 성공적으로 보냈습니다.')
        else:
            print(response)


def make_multimsg(msg_dict):
    multi = MIMEMultipart(_subtype='mixed')

    for key, value in msg_dict.items():
        # 각 타입에 적절한 MIMExxx() 함수를 호출하여 msg 객체를 생성한다.
        if key == 'text':
            with open(value['filename'], encoding='utf-8') as fp:
                msg = MIMEText(fp.read(), _subtype=value['subtype'])
        elif key == 'image':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEImage(fp.read(), _subtype=value['subtype'])
        elif key == 'audio':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEAudio(fp.read(), _subtype=value['subtype'])
        else:
            with open(value['filename'], 'rb') as fp:
                msg = MIMEBase(value['maintype'], _subtype=value['subtype'])
                msg.set_payload(fp.read())
                encoders.encode_base64(msg)
        # 파일 이름을 첨부파일 제목으로 추가
        msg.add_header('Content-Disposition', 'attachment',
                       filename=os.path.basename(value['filename']))
        # 첨부파일 추가
        multi.attach(msg)

    return multi






#################################
## 함수 정의
#################################
def get_stock_code():
    # 종목코드 다운로드
    stock_code = pd.read_html('http://kind.krx.co.kr/corpgeneral/corpList.do?method=download', header=0)[0]

    # 필요없는 column들은 제외
    stock_code = stock_code[['회사명', '종목코드']]

    # 한글 컬럼명을 영어로 변경
    stock_code = stock_code.rename(columns={'회사명': 'company', '종목코드': 'code'})

    # 종목코드가 6자리이기 때문에 6자리를 맞춰주기 위해 설정해줌
    stock_code.code = stock_code.code.map('{:06d}'.format)

    return stock_code

def get_stock(code):
    df = pd.DataFrame()
    for page in range(1,21):
        # 일별 시세 url
        url = 'http://finance.naver.com/item/sise_day.nhn?code={code}'.format(code=code)
        url = '{url}&page={page}'.format(url=url, page=page)
        print(url)
        header = {'User-Agent':'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/106.0.0.0 Safari/537.36'} ## 참고) 7.3.1 User-Agent 확인하기 
        res = requests.get(url, headers=header)
        current_df = pd.read_html(res.text, header=0)[0]
        df = df.append(current_df, ignore_index=True)

    return df

def clean_data(df):
    # df.dropna()를 이용해 결측값 있는 행 제거
    df = df.dropna()

    # 한글로 된 컬럼명을 영어로 바꿔줌
    df = df.rename(columns= {'날짜': 'date', '종가': 'close', '전일비': 'diff','시가': 'open', '고가': 'high', '저가': 'low','거래량': 'volume'})
    # 데이터의 타입을 int형으로 바꿔줌
    df[['close', 'diff', 'open', 'high', 'low', 'volume']] =df[['close', 'diff', 'open', 'high', 'low', 'volume']].astype(int)

    # 컬럼명 'date'의 타입을 date로 바꿔줌
    df['date'] = pd.to_datetime(df['date'])

    # 일자(date)를 기준으로 오름차순 정렬
    df = df.sort_values(by=['date'], ascending=True)

    return df

#################################
## 함수 호출
#################################
# 종목코드 가져오기
company='삼성전자'
stock_code = get_stock_code()

# 일별 시세 가져오기
code = stock_code[stock_code.company==company].code.values[0].strip() ## strip() : 공백제거
df = get_stock(code)

# 일별 시세 클린징
df = clean_data(df)
#print(df)


#################################
## 차트 그리기
#################################
plt.figure(figsize=(10,4))
plt.plot(df['date'], df['close'])
plt.xlabel('date')
plt.ylabel('close')

#################################
## 차트 저장 및 출력하기
#################################
chart_fname = os.path.join("res/stock_report" ,'{company}_chart.png'.format(company=company))
plt.savefig(chart_fname)
plt.show()

#################################
## 일별 시세 그리기
#################################
plt.figure(figsize=(15,4))
ax = plt.subplot(111, frame_on=False) # no visible frame
ax.xaxis.set_visible(False) # hide the x axis
ax.yaxis.set_visible(False) # hide the y axis
df = df.sort_values(by=['date'], ascending=False)
table(ax, df.head(10), loc='center', cellLoc = 'center', rowLoc = 'center')# where df is your data frame

#################################
## 일별 시세 저장하기
#################################
table_fname = os.path.join("res/stock_report" ,'{company}_table.png'.format(company=company))
plt.savefig(table_fname)


#################################
## 파워포인트 객체 선언
#################################
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation() # 파워포인트 객체 선언

#################################
## 제목 슬라이드 추가
#################################
title_slide_layout = prs.slide_layouts[0] # 0 : 제목 슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout) # 제목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title # 제목
title.text = "주식 보고서" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목 상자는 placeholders[0], 부제목 상자는 [1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

#################################
## 차트 및 테이블 슬라이드 추가
#################################
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company,close=df.iloc[0]['close'])
print(shapes.title.text)

# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
# width, hegith가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)

pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element) # 해당 요소를 뒤로 보냅니다.

#################################
## 보고서 저장
#################################
ppt_fname = os.path.join("res/stock_report" ,'stock_report.pptx')
prs.save(ppt_fname)



smtp_info = dict({"smtp_server" : "smtp.naver.com", # SMTP 서버 주소
                  "smtp_user_id" : "yclee7777@naver.com",
                  "smtp_user_pw" : "qustls&14NV",
                  "smtp_port" : 587}) # SMTP 서버 포트

msg_dict = {
    #'text' : {'maintype' : 'text', 'subtype' :'plain', 'filename' : 'res/email_sending/test1.txt'}, # 텍스트 첨부파일
    #'image' : {'maintype' : 'image', 'subtype' :'jpg', 'filename' : 'res/email_sending/test.jpg' }, # 이미지 첨부파일
    #'audio' : {'maintype' : 'audio', 'subtype' :'mp3', 'filename' : 'res/email_sending/test.mp3' }, # 오디오 첨부파일
    #'video' : {'maintype' : 'video', 'subtype' :'mp4', 'filename' : 'res/email_sending/test.mp4' }, # 비디오 첨부파일
    'application' : {'maintype' : 'application', 'subtype' : 'octect-stream', 'filename' : 'res/email_sending/test.pdf'} # 그 외 첨부파일
}


# # 메일 내용 작성
# title = "기본 이메일 입니다.[나에게]"
# content = "메일 내용입니다. 파이선 메일 송신 테스트"
# sender = smtp_info['smtp_user_id'] # 송신자(sender) 메일 계정
# receiver = "yclee7777@naver.com"

# # 메일 객체 생성 : 메시지 내용에는 한글이 들어가기 때문에 한글을 지원하는 문자 체계인 UTF-8을 명시해줍니다.
# msg = MIMEText(_text = content, _charset = "utf-8") # 메일 내용

# msg['Subject'] = title # 메일 제목
# msg['From'] = sender # 송신자
# msg['To'] = receiver # 수신자

# send_email(smtp_info, msg )


#####################
# 메일 내용 작성
#####################
title = '({date}). 주식 보고서 분석 자료 입니다.'.format(date = today)
content = "주식 보고서 분서 자료 입니다."
sender = smtp_info['smtp_user_id'] # 송신자(sender) 메일 계정
receiver = "yclee7777@naver.com"

# 메일 내용
msg = MIMEText(_text = content, _charset = "utf-8")


# 첨부파일 추가
msg_dict['application']['filename'] = ppt_fname
multi = make_multimsg(msg_dict)
multi['Subject'] = title
multi['From'] = sender
multi['To'] = receiver
multi.attach(msg)

# 첨부파일이 추가된 이메일 전송
send_email(smtp_info, multi )

