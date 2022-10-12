from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해


img_path = "res/powerpoint_handling/lee1.jpg"

prs = Presentation() # 파워포인트 객체 선언

blank_slide_layout = prs.slide_layouts[6] # 6 : 제목 없는 빈 슬라이드에 해당
slide = prs.slides.add_slide(blank_slide_layout) # 슬라이드 추가

left = top = Inches(1)
width = height = Inches(1)
# width, hegith가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(img_path, left, top, width=width,height=height)

left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, width=width,height=height)


shapes = slide.shapes
rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table



# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


# 저장할 파일명
PPT_FILENAME = "res/powerpoint_handling/test7.pptx"


# 저장
prs.save(PPT_FILENAME)