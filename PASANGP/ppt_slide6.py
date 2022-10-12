from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해


prs = Presentation() # 파워포인트 객체 선언

title_only_slide_layout = prs.slide_layouts[5]         # 5 제목 슬라이드 
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

title_shape = slide.placeholders[0] 
title_shape.text = 'Adding a Table'

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

# write column headings
table.cell(0, 0).text = 'Foo'
table.cell(0, 1).text = 'Bar'

# write body cells
table.cell(1, 0).text = 'Baz'
table.cell(1, 1).text = 'Qux'


# 카카오 토큰을 저장할 파일명
PPT_FILENAME = "res/powerpoint_handling/test6.pptx"


# 저장
prs.save(PPT_FILENAME)