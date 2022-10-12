from pptx import Presentation # 라이브러리
from pptx.util import Inches # 사진, 표등을 그리기 위해

prs = Presentation() # 파워포인트 객체 선언

title_slide_layout = prs.slide_layouts[0] # 0 : 제목 슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout) # 슬라이드 추가

# 제목 - 제목에 값넣기
title = slide.placeholders[0] # 제목
title.text = "Hello, World!" # 제목에 값 넣기

# 부제목
subtitle = slide.placeholders[1] # 제목 상자는 placeholders[0], 부제목 상자는 [1]
subtitle.text = "python-pptx was here!" 

# 카카오 토큰을 저장할 파일명
PPT_FILENAME = "res/powerpoint_handling/test1.pptx"


# 저장
prs.save(PPT_FILENAME)